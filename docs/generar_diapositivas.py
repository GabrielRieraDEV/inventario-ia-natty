"""Genera la presentación de la defensa del TEG (PowerPoint editable).

Sigue la estructura oficial exigida por la universidad y usa el contenido
real del TEG (docs/TEG.tex). Salida: docs/Defensa_TEG.pptx

Uso:  python generar_diapositivas.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Paleta institucional UGMA (azul y dorado) — tono sobrio ----
NAVY = RGBColor(0x12, 0x3A, 0x6B)     # azul UGMA (profundo, institucional)
ACCENT = RGBColor(0xC0, 0x9A, 0x4A)   # dorado apagado (sobrio, no estridente)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x29, 0x33)     # gris muy oscuro para el texto
GRAY = RGBColor(0x53, 0x60, 0x6E)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

_n = [0]  # contador de diapositivas de contenido


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb


def _set(p, text, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False):
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"


def base(slide):
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.15))
    _fill(bar, NAVY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.15), SW, Pt(4))
    _fill(line, ACCENT)


def titulo(slide, texto, tag=None):
    base(slide)
    tb = _box(slide, Inches(0.6), Inches(0.18), Inches(12.1), Inches(0.85))
    tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(tb.text_frame.paragraphs[0], texto, 26, WHITE, bold=True)
    if tag:
        tg = _box(slide, Inches(0.62), Inches(0.86), Inches(12), Inches(0.3))
        _set(tg.text_frame.paragraphs[0], tag, 12, ACCENT, bold=True)


def pie(slide):
    _n[0] += 1
    f = _box(slide, Inches(0.5), Inches(7.0), Inches(9), Inches(0.4))
    _set(f.text_frame.paragraphs[0], "Defensa TEG · Gabriel Riera · UGMA Núcleo Bolívar", 9, GRAY)
    pg = _box(slide, Inches(12.3), Inches(7.0), Inches(0.7), Inches(0.4))
    _set(pg.text_frame.paragraphs[0], str(_n[0]), 11, NAVY, bold=True, align=PP_ALIGN.RIGHT)


def vinetas(slide, items, top=1.55, left=0.8, width=11.8, height=5.2, size=18):
    """items: lista de (texto, nivel)."""
    tb = _box(slide, Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        p.line_spacing = 1.05
        if lvl == 0:
            _set(p, "•  " + txt, size, DARK, bold=False)
        elif lvl == 1:
            _set(p, "–  " + txt, size - 2, GRAY)
            p.level = 1
        else:  # título de bloque
            _set(p, txt, size, NAVY, bold=True)
            p.space_before = Pt(6)
    return tb


def placeholder(slide, l, t, w, h, label):
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = NAVY; box.line.width = Pt(1.25); box.line.dash_style = 2
    tf = box.text_frame; tf.word_wrap = True
    _set(tf.paragraphs[0], "🖼  " + label, 12, GRAY, align=PP_ALIGN.CENTER, italic=True)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def contenido(titulo_txt, items, tag=None, size=18):
    s = prs.slides.add_slide(BLANK)
    titulo(s, titulo_txt, tag)
    vinetas(s, items, size=size)
    pie(s)
    return s


# =====================================================================
# 1. PORTADA
# =====================================================================
s = prs.slides.add_slide(BLANK)
from pptx.enum.shapes import MSO_SHAPE
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
_fill(bg, NAVY)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.55), SW, Pt(3))
_fill(band, ACCENT)

t = _box(s, Inches(1), Inches(0.5), Inches(11.3), Inches(1.5))
_set(t.text_frame.paragraphs[0], "Universidad Nororiental Privada Gran Mariscal de Ayacucho", 16, WHITE, bold=True, align=PP_ALIGN.CENTER)
for extra in ["Escuela de Ingeniería en Informática — Núcleo Bolívar",
              "Área de Conocimiento: Ingeniería de Software"]:
    p = t.text_frame.add_paragraph(); _set(p, extra, 12, ACCENT, align=PP_ALIGN.CENTER)

# logo placeholder
lg = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.07), Inches(1.9), Inches(1.2), Inches(1.2))
lg.fill.solid(); lg.fill.fore_color.rgb = WHITE; lg.line.fill.background()
_set(lg.text_frame.paragraphs[0], "LOGO\nUGMA", 10, NAVY, bold=True, align=PP_ALIGN.CENTER)
lg.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

ti = _box(s, Inches(1), Inches(3.25), Inches(11.3), Inches(1.2))
_set(ti.text_frame.paragraphs[0],
     "Sistema de Digitalización y Registro de Productos con Inteligencia Artificial para el Control de Inventarios",
     24, WHITE, bold=True, align=PP_ALIGN.CENTER)
p = ti.text_frame.add_paragraph()
_set(p, "Bodegón Pa' Donde Natty C.A. — Ciudad Bolívar, estado Bolívar", 15, ACCENT, align=PP_ALIGN.CENTER)

meta = _box(s, Inches(1), Inches(4.8), Inches(11.3), Inches(2.2))
datos = [
    ("Trabajo Especial de Grado para optar al título de Ingeniero en Informática", 13, WHITE, False),
    ("", 6, WHITE, False),
    ("Autor: Gabriel Riera — C.I.: V-31.153.478", 15, WHITE, True),
    ("Tutor Académico: Ing. Carlos Moreno", 13, WHITE, False),
    ("Ciudad Bolívar, 2026", 12, ACCENT, False),
]
for i, (txt, sz, col, bd) in enumerate(datos):
    p = meta.text_frame.paragraphs[0] if i == 0 else meta.text_frame.add_paragraph()
    _set(p, txt, sz, col, bold=bd, align=PP_ALIGN.CENTER)

# =====================================================================
# 2. CONTENIDO (agenda)
# =====================================================================
contenido("Contenido", [
    ("Introducción", 0),
    ("Planteamiento del problema", 0),
    ("Objetivos de la investigación", 0),
    ("Justificación y alcance", 0),
    ("Marco teórico", 0),
    ("Metodología", 0),
    ("Resultados (desarrollo de los objetivos específicos)", 0),
    ("Conclusiones y recomendaciones", 0),
], size=20)

# =====================================================================
# 3. INTRODUCCIÓN
# =====================================================================
contenido("Introducción", [
    ("La gestión eficiente del inventario es clave para la competitividad de las empresas comerciales.", 0),
    ("La transformación digital y la Inteligencia Artificial permiten optimizar procesos y decisiones en tiempo real.", 0),
    ("En Venezuela, muchas pymes y bodegas aún dependen de registros manuales o desactualizados.", 0),
    ("El Bodegón Pa' Donde Natty C.A. (Ciudad Bolívar) controla su inventario de forma manual: conteos en papel y hojas de cálculo no integradas.", 0),
    ("Este trabajo propone un prototipo que digitaliza el registro de productos usando un modelo de IA multimodal que reconoce el producto a partir de una foto.", 0),
])

# =====================================================================
# 4. PLANTEAMIENTO DEL PROBLEMA
# =====================================================================
contenido("Planteamiento del Problema", [
    ("El control de existencias es manual: conteos físicos, apuntes en papel y hojas de cálculo no sincronizadas.", 0),
    ("Errores recurrentes detectados:", 2),
    ("Discrepancia entre el inventario físico y el registrado.", 1),
    ("Desabastecimiento de productos de alta rotación.", 1),
    ("Pérdidas por vencimiento no detectado a tiempo.", 1),
    ("Duplicidad de registros y retardo en la reposición.", 1),
    ("Consecuencia: imposibilidad de una gestión dinámica y confiable, con pérdidas económicas evitables.", 0),
    ("Pregunta central: ¿cómo lograr un registro de inventario exacto y en tiempo real apoyado en IA?", 0),
])

# =====================================================================
# 5. OBJETIVOS
# =====================================================================
contenido("Objetivos de la Investigación", [
    ("Objetivo General", 2),
    ("Desarrollar un prototipo funcional de un sistema de digitalización y registro de productos con IA para el control de inventarios en el Bodegón Pa' Donde Natty C.A.", 0),
    ("Objetivos Específicos", 2),
    ("1. Diagnosticar las necesidades funcionales y técnicas del proceso actual de inventario.", 0),
    ("2. Diseñar la arquitectura, el modelo de datos y la interfaz (consolidados en una SRS).", 0),
    ("3. Desarrollar el prototipo integrando el reconocimiento mediante la API de un modelo multimodal.", 0),
    ("4. Validar el funcionamiento del prototipo mediante pruebas en un entorno controlado.", 0),
], size=17)

# =====================================================================
# 6. JUSTIFICACIÓN Y ALCANCE
# =====================================================================
contenido("Justificación y Alcance", [
    ("Justificación", 2),
    ("Tecnológica: moderniza la operación integrando IA accesible (sin entrenar un modelo propio).", 0),
    ("Empresarial: ataca la causa raíz (falta de registro centralizado y oportuno) y reduce pérdidas.", 0),
    ("Académica/social: modelo replicable para otras pymes del sector comercial del estado Bolívar.", 0),
    ("Alcance", 2),
    ("Prototipo funcional con módulos de reconocimiento IA, registro de movimientos, alertas, consulta, reportes, catálogo, autenticación, dashboard y respaldo.", 0),
    ("Validado en entorno controlado; desplegado en capa gratuita (no implantación definitiva en la operación real).", 0),
], size=16)

# =====================================================================
# 7. MARCO TEÓRICO
# =====================================================================
contenido("Marco Teórico", [
    ("Antecedente más relevante", 2),
    ("Wong Portillo, De la Torre y Reyes (2023, UPC, Perú): «Aplicación de reconocimiento de imágenes para productos en el sector retail».", 0),
    ("Demostraron la viabilidad de reconocer productos por imagen (modelos YOLO). Aporte: este trabajo logra el mismo fin con un modelo multimodal vía API, sin entrenar un modelo propio.", 1),
    ("Términos clave", 2),
    ("Control de inventarios: supervisión y registro de existencias para garantizar disponibilidad minimizando costos.", 0),
    ("Inteligencia Artificial / aprendizaje automático: sistemas que aprenden de datos y reconocen patrones.", 0),
    ("Modelo de lenguaje multimodal: IA que relaciona texto e imagen y responde en lenguaje natural.", 0),
    ("Reconocimiento de imágenes: identificación del contenido de una imagen (aquí, el producto).", 0),
    ("API e ingeniería de prompts: consumo de un servicio de IA externo guiado por una instrucción bien diseñada.", 0),
], size=15)

# =====================================================================
# 8. METODOLOGÍA (1)
# =====================================================================
contenido("Metodología (1/3)", [
    ("Tipo y diseño de la investigación", 2),
    ("Tipo: Descriptiva, Aplicada y Proyectiva (propone, desarrolla y valida una solución).", 0),
    ("Diseño: de Campo (datos tomados en la empresa) con componente Documental; no experimental.", 0),
    ("Población y muestra", 2),
    ("Población: 5 trabajadores del bodegón vinculados al inventario.", 0),
    ("Muestra: censo poblacional (los 5 trabajadores), por ser población finita y reducida.", 0),
    ("Técnicas de recolección de datos", 2),
    ("Encuesta (cuestionario), observación directa y entrevista, análisis documental (SRS / MTR), y pruebas de software (lista de cotejo).", 0),
], tag="OBJETIVOS 1 A 4", size=15)

# =====================================================================
# 9. METODOLOGÍA (2) — análisis de datos
# =====================================================================
contenido("Metodología (2/3) — Análisis de Datos", [
    ("Diagnóstico (Obj. 1): estadística descriptiva (frecuencias y porcentajes) sobre el cuestionario, triangulada con observación y entrevista.", 0),
    ("Diseño (Obj. 2): análisis documental de la SRS y verificación con la Matriz de Trazabilidad de Requisitos (cobertura, acoplamiento, normalización 3FN).", 0),
    ("Desarrollo (Obj. 3): observación sistemática del porcentaje de módulos construidos y de la integración correcta de la API de IA.", 0),
    ("Validación (Obj. 4): pruebas de caja negra y lista de cotejo (Cumple / No cumple); medición de exactitud y tiempo de respuesta del reconocimiento.", 0),
], size=17)

# =====================================================================
# 10. METODOLOGÍA (3) — operacionalización
# =====================================================================
contenido("Metodología (3/3) — Operacionalización de la Variable", [
    ("Variable: prototipo de sistema de digitalización y registro de productos con IA para el control de inventarios.", 0),
    ("Obj. 1 — Dim. Operativa: errores de registro, tiempo por ciclo, necesidades priorizadas → Encuesta/Observación/Entrevista.", 0),
    ("Obj. 2 — Dim. Técnica (diseño): acoplamiento, normalización 3FN, cobertura de la SRS → Análisis documental / MTR.", 0),
    ("Obj. 3 — Dim. Funcional: % de módulos, integración de la API, efectividad del reconocimiento → Kanban / Observación.", 0),
    ("Obj. 4 — Dim. Validación: tasa de cumplimiento, exactitud ≥ 90 %, tiempo ≤ 3 s → Pruebas de caja negra / Lista de cotejo.", 0),
], size=16)

# =====================================================================
# 11+. RESULTADOS (objetivos específicos)
# =====================================================================
contenido("Resultados — Obj. 1: Diagnóstico (1/2)", [
    ("Proceso actual: control manual, sin registro centralizado ni actualización en tiempo real.", 0),
    ("Frecuencia de errores (censo de 5 trabajadores):", 2),
    ("Discrepancia físico vs. registrado: 100 %.", 1),
    ("Desabastecimiento de alta rotación: 80 %  ·  Pérdidas por vencimiento: 80 %.", 1),
    ("Duplicidad de registros: 60 %  ·  Retardo en reposición: 60 %.", 1),
    ("Tiempo del ciclo de inventario manual: 4 a 6 horas, con la operación parcialmente interrumpida.", 0),
], tag="OBJETIVO ESPECÍFICO 1", size=17)

contenido("Resultados — Obj. 1: Requerimientos (2/2)", [
    ("El diagnóstico se consolidó en requerimientos formales y verificables:", 0),
    ("9 Requerimientos Funcionales (RF-01 a RF-09)", 2),
    ("Reconocimiento IA por foto (QR), movimientos, alertas, consulta, reportes, catálogo, autenticación, dashboard y respaldo.", 1),
    ("7 Requerimientos Técnicos (RT-01 a RT-07)", 2),
    ("Backend Python/FastAPI + IA por API; BD PostgreSQL/Supabase; reconocimiento ≤ 3 s; exactitud ≥ 90 %; web + QR; respaldo; despliegue en capa gratuita.", 1),
], tag="OBJETIVO ESPECÍFICO 1", size=16)

s = contenido("Resultados — Obj. 2: Diseño", [
    ("Arquitectura modular de 3 capas con bajo acoplamiento:", 0),
    ("Presentación (Next.js, PC + móvil) → Lógica de negocio (FastAPI + cliente de IA) → Persistencia (PostgreSQL/Supabase).", 1),
    ("Modelo de datos normalizado a 3FN; el stock es un dato derivado y auditable de los movimientos.", 0),
    ("Matriz de Trazabilidad (MTR): cobertura del 100 % (9/9 RF), sin requerimientos huérfanos.", 0),
], tag="OBJETIVO ESPECÍFICO 2", size=16)
placeholder(s, 8.7, 1.6, 4.2, 2.5, "Insertar: diagrama de\narquitectura / ERD")

s = contenido("Resultados — Obj. 3: Desarrollo del Prototipo", [
    ("Construidos los 9 módulos del diseño (metodología Kanban).", 0),
    ("Stack: Python/FastAPI · Next.js · PostgreSQL/Supabase · Google Gemini (multimodal).", 0),
    ("Flujo estrella (RF-01): la PC genera un QR → el teléfono toma la foto → la IA reconoce el producto → se registra el movimiento y se refleja en la PC.", 0),
    ("Integración por API y prompt estructurado: respuesta directamente procesable por la lógica de inventario.", 0),
], tag="OBJETIVO ESPECÍFICO 3", size=16)
placeholder(s, 8.9, 4.2, 4.0, 2.4, "Insertar capturas:\nQR + captura móvil")

s = contenido("Resultados — Obj. 3: Reconocimiento por IA", [
    ("Se delega la visión en un modelo multimodal de propósito general (no se entrena un modelo propio).", 0),
    ("Desempeño medido sobre 100 imágenes de 20 productos distintos:", 2),
    ("Exactitud del reconocimiento: 93 %  (supera el umbral RT-04 ≥ 90 %).", 1),
    ("Tiempo de respuesta promedio: 2,4 s  (cumple RT-03 ≤ 3 s).", 1),
    ("Decisión de diseño validada: menor costo y complejidad que entrenar un modelo propio (vs. Wong Portillo et al., 2023).", 0),
], tag="OBJETIVO ESPECÍFICO 3", size=16)
placeholder(s, 9.0, 4.3, 3.9, 2.3, "Insertar captura:\npantalla de reconocimiento")

contenido("Resultados — Obj. 4: Validación", [
    ("Pruebas de caja negra + lista de cotejo sobre cada requerimiento funcional.", 0),
    ("Tasa de cumplimiento de requisitos funcionales: 100 % (9/9 «Cumple»).", 0),
    ("Exactitud del reconocimiento: 93 %  ·  Tiempo de respuesta: 2,4 s.", 0),
    ("Registro automatizado sin discrepancias frente a los casos verificados manualmente (margen de error nulo).", 0),
    ("Cada error diagnosticado tiene hoy un requerimiento que lo combate, un componente que lo implementa y una prueba que lo confirma.", 0),
], tag="OBJETIVO ESPECÍFICO 4", size=16)

contenido("Resultados — Factibilidad", [
    ("Técnica: herramientas de código abierto y despliegue en capa gratuita (Vercel, Render/HF, Supabase); desempeño verificado (93 % / 2,4 s).", 0),
    ("Operativa: interfaz sencilla y consistente; sustituye un proceso de 4–6 h por un registro asistido por imagen (curva de aprendizaje corta).", 0),
    ("Económica: inversión estimada ≈ 1.150 USD, recuperable a corto plazo frente al ahorro por menos pérdidas y errores.", 0),
], tag="COMPLEMENTARIO", size=17)

# =====================================================================
# CONCLUSIONES
# =====================================================================
contenido("Conclusiones", [
    ("Obj. 1: el diagnóstico (triangulado) confirmó la discrepancia del 100 % y ciclos de 4–6 h; se consolidó en 9 RF y 7 RT verificables.", 0),
    ("Obj. 2: la arquitectura de 3 capas (bajo acoplamiento), el modelo 3FN y la MTR al 100 % cubren íntegramente las necesidades.", 0),
    ("Obj. 3: se desarrollaron los 9 módulos; el reconocimiento por API alcanzó 93 % de exactitud y 2,4 s.", 0),
    ("Obj. 4: validación con 100 % de cumplimiento y margen de error nulo en el registro automatizado.", 0),
    ("General: integrar un modelo multimodal vía API es una alternativa viable, accesible y eficaz para digitalizar el inventario de pymes comerciales.", 0),
], size=16)

# =====================================================================
# RECOMENDACIONES
# =====================================================================
contenido("Recomendaciones", [
    ("A la empresa", 2),
    ("Implementar de forma gradual (empezando por una categoría) y capacitar al personal en la captura por QR.", 0),
    ("Política de respaldos y auditoría periódica físico vs. registrado.", 0),
    ("Técnicas / continuidad", 2),
    ("Refinar el prompt con los casos no reconocidos (7 % restante) y añadir caché de productos frecuentes.", 0),
    ("Incorporar un módulo de auditoría (usuario, fecha y hora de cada movimiento).", 0),
    ("Investigaciones futuras", 2),
    ("Predicción de demanda con aprendizaje automático; facturación/contabilidad; replicar el estudio en otras pymes.", 0),
], size=15)

# =====================================================================
# MENSAJE MOTIVADOR
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
_fill(bg, NAVY)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.05), SW, Pt(3))
_fill(band, ACCENT)
m = _box(s, Inches(1.2), Inches(2.4), Inches(10.9), Inches(2.6))
m.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
_set(m.text_frame.paragraphs[0],
     "«La tecnología no reemplaza el esfuerzo humano: lo potencia.",
     26, WHITE, bold=True, align=PP_ALIGN.CENTER)
p = m.text_frame.add_paragraph()
_set(p, "Digitalizar es dar el primer paso hacia un negocio más justo, preciso y competitivo.»",
     20, ACCENT, bold=True, align=PP_ALIGN.CENTER)
p = m.text_frame.add_paragraph()
_set(p, "\n¡Gracias por su atención!", 18, WHITE, align=PP_ALIGN.CENTER)

prs.save("Defensa_TEG.pptx")
print(f"Listo: Defensa_TEG.pptx con {len(prs.slides.__iter__.__self__._sldIdLst)} diapositivas.")
